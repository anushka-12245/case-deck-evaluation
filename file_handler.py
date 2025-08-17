import pdfplumber
from pptx import Presentation
from PIL import Image
import io
import tempfile
import pytesseract

def extract_text_from_file(filename, content):
    if filename.endswith(".pdf"):
        return extract_from_pdf(content)
    elif filename.endswith((".ppt", ".pptx")):
        return extract_from_ppt(content)
    elif filename.endswith((".jpg", ".jpeg", ".png")):
        return extract_from_image(content)
    return "Unsupported file format."

def extract_from_pdf(content):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(content)
        tmp.flush()
        with pdfplumber.open(tmp.name) as pdf:
            text = "\n".join(page.extract_text() or "" for page in pdf.pages)
    return text.strip()

def extract_from_ppt(content):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(content)
        tmp.flush()
        prs = Presentation(tmp.name)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    return "\n".join(text).strip()

def extract_from_image(content):
    image = Image.open(io.BytesIO(content))
    return pytesseract.image_to_string(image).strip()
